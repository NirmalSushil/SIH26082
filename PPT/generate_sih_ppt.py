import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def generate_perfect_sih_presentation():
    template_path = 'PPT/SIH2026-IDEA-Presentation-Format.pptx'
    output_path = 'PPT/SIH2026-IDEA-Presentation-Xtream-09.pptx'
    
    prs = Presentation(template_path)
    
    # -------------------------------------------------------------
    # BRAND COLORS & TYPOGRAPHY
    # -------------------------------------------------------------
    COLOR_PRIMARY = RGBColor(16, 85, 50)      # Forest / Emerald Green (Clean & Green Theme)
    COLOR_ACCENT = RGBColor(5, 150, 105)      # Vibrant Emerald
    COLOR_HEADER = RGBColor(15, 23, 42)       # Slate 900
    COLOR_BODY = RGBColor(30, 41, 59)         # Slate 800
    COLOR_HIGHLIGHT = RGBColor(2, 132, 199)   # Sky Blue
    
    # -------------------------------------------------------------
    # SLIDE 1: TITLE PAGE
    # -------------------------------------------------------------
    slide1 = prs.slides[0]
    
    # Clean up Subtitle placeholder if needed
    for shape in slide1.shapes:
        if shape.name == "Subtitle 3" and shape.has_text_frame:
            shape.text_frame.clear()
            p = shape.text_frame.paragraphs[0]
            p.text = "SMART INDIA HACKATHON 2026"
            p.font.size = Pt(26)
            p.font.bold = True
            p.font.color.rgb = COLOR_PRIMARY
            p.alignment = PP_ALIGN.LEFT

        if shape.name == "Title 7" and shape.has_text_frame:
            shape.text_frame.clear()
            
        if shape.name == "TextBox 9" and shape.has_text_frame:
            # Position cleanly
            shape.left = Inches(0.67)
            shape.top = Inches(1.85)
            shape.width = Inches(6.80)
            shape.height = Inches(5.10)
            
            tf = shape.text_frame
            tf.word_wrap = True
            tf.clear()
            
            fields = [
                ("Problem Statement ID:", " SIH26082", True),
                ("Problem Statement Title:", " Air Pollution - Weather Coupled Forecast", True),
                ("Theme:", " Clean & Green Technology", True),
                ("PS Category:", " Software", False),
                ("Team ID:", " 09", True),
                ("Team Name:", " Xtream", True),
            ]
            
            for idx, (label, value, is_primary) in enumerate(fields):
                p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
                p.space_after = Pt(14)
                
                run1 = p.add_run()
                run1.text = label
                run1.font.bold = True
                run1.font.size = Pt(16)
                run1.font.color.rgb = COLOR_HEADER
                
                run2 = p.add_run()
                run2.text = value
                run2.font.bold = True
                run2.font.size = Pt(17)
                if "Clean & Green" in value or "Xtream" in value or "SIH26082" in value:
                    run2.font.color.rgb = COLOR_PRIMARY
                else:
                    run2.font.color.rgb = COLOR_BODY

    # -------------------------------------------------------------
    # SLIDES 2 to 6: CONTENT SLIDES
    # -------------------------------------------------------------
    content_slides_data = [
        # SLIDE 2
        {
            "slide_num": 1,
            "title": "IDEA: ATMOCAST - Weather-Coupled Air Quality Forecast",
            "sections": [
                {
                    "heading": "1. Proposed Solution:",
                    "points": [
                        ("Physics-Informed Atmospheric Coupling:", " Integrates real-time chemical air quality metrics with meteorological dynamics to generate high-resolution, short-term coupled forecasts."),
                        ("Dynamic Multi-Factor Modeling:", " Accurately simulates how wind shear, boundary layer thermal inversions, rain scavenging, and barometric pressure modify ground-level pollution concentrations.")
                    ]
                },
                {
                    "heading": "2. How It Addresses the Problem:",
                    "points": [
                        ("Bridges Static Observation Gaps:", " Replaces lagging ground sensor measurements with proactive microclimate-informed forecasts."),
                        ("Anticipates Smog Trapping & Clearing:", " Detects nocturnal inversion accumulation and daytime convective dispersion windows before they occur.")
                    ]
                },
                {
                    "heading": "3. Innovation & Uniqueness:",
                    "points": [
                        ("Transparent Explainability:", " Provides an exact mathematical delta breakdown (+/- AQI) attributing numerical points to wind, rain, and inversion."),
                        ("Zero-Cost Keyless Architecture:", " Operates entirely on public Open-Meteo APIs with a deterministic physical fallback model for 100% uptime.")
                    ]
                }
            ]
        },
        # SLIDE 3
        {
            "slide_num": 2,
            "title": "TECHNICAL APPROACH & IMPLEMENTATION",
            "sections": [
                {
                    "heading": "1. Technology Stack & Frameworks:",
                    "points": [
                        ("Core Architecture:", " Next.js 16 (App Router, React 19) + TypeScript 5 (Strict Mode) for high-performance edge compute."),
                        ("UI & Visualization:", " Tailwind CSS v4, enterprise design tokens (Light/Dark mode), and Recharts 3.10 multi-axis interactive composed charts."),
                        ("Data Integration:", " Open-Meteo Forecast, Air Quality, and Geocoding APIs coupled with a built-in synthetic atmospheric model.")
                    ]
                },
                {
                    "heading": "2. Methodology & Coupling Pipeline:",
                    "points": [
                        ("Atmospheric Ingestion:", " Streams 48-hour hourly weather (wind, gusts, temp, humidity, pressure, rain) & pollutants (PM2.5, PM10, O3, NO2, SO2, CO)."),
                        ("Coupling Heuristic Engine:", " Calculates: Coupled AQI = Raw AQI + Delta_wind + Delta_stability + Delta_precip + Delta_photochem + Delta_pressure."),
                        ("Actionable Analytics Deck:", " Real-time Hero AQI card, 24h timeline reel, factor diagnostics, and optimal outdoor activity recommendations.")
                    ]
                }
            ]
        },
        # SLIDE 4
        {
            "slide_num": 3,
            "title": "FEASIBILITY, VIABILITY & RISK MITIGATION",
            "sections": [
                {
                    "heading": "1. Feasibility & Operational Viability:",
                    "points": [
                        ("Zero API Overhead & High Availability:", " Free, keyless public infrastructure eliminates subscription costs and quota restrictions."),
                        ("Ultra-Low Latency Execution:", " Mathematical coupling heuristics compute complete 48-hour trajectories in <50ms per query."),
                        ("Cross-Platform Responsive Access:", " Accessible on all devices with 1-click location bookmarking and GPS integration.")
                    ]
                },
                {
                    "heading": "2. Potential Challenges & Risks:",
                    "points": [
                        ("External Network Latency / API Downtime:", " Risk of third-party server slowdowns or connection timeouts."),
                        ("Microclimatic Urban Variations:", " Localized urban street canyons differing from macro meteorological stations.")
                    ]
                },
                {
                    "heading": "3. Risk Mitigation Strategies:",
                    "points": [
                        ("Synthetic Atmospheric Fallback:", " Automatically activates physical diurnal climate algorithms during network timeouts to guarantee zero-downtime."),
                        ("Pre-Indexed Global City Database:", " Instant local geocoding for 30+ major world hubs without network round-trips.")
                    ]
                }
            ]
        },
        # SLIDE 5
        {
            "slide_num": 4,
            "title": "IMPACT, BENEFITS & SUSTAINABILITY (CLEAN & GREEN)",
            "sections": [
                {
                    "heading": "1. Target Audience Impact:",
                    "points": [
                        ("Citizens & Commuters:", " Real-time visibility into peak smog windows and cleanest outdoor travel/exercise hours."),
                        ("Sensitive Demographics:", " Proactive alerts for asthmatics, children, and elderly citizens to prevent acute respiratory episodes."),
                        ("Municipal Authorities:", " Data-driven triggers for anti-smog guns, street misting, and construction dust regulations.")
                    ]
                },
                {
                    "heading": "2. Multi-Dimensional Benefits:",
                    "points": [
                        ("Social & Healthcare:", " Substantially lowers emergency hospital visits through forward-looking smog warnings."),
                        ("Economic:", " Reduces healthcare expenditures and optimizes municipal pollution-mitigation budgets."),
                        ("Environmental (Clean & Green):", " Aligns with National Clean Air Programme (NCAP) and UN SDGs (SDG 3: Good Health & SDG 11: Sustainable Cities).")
                    ]
                }
            ]
        },
        # SLIDE 6
        {
            "slide_num": 5,
            "title": "RESEARCH, REFERENCES & STANDARDS",
            "sections": [
                {
                    "heading": "1. Environmental Standards & Regulatory Frameworks:",
                    "points": [
                        ("US EPA Air Quality Index Standards:", " Technical Assistance Document for Daily Air Quality Reporting (40 CFR Part 58)."),
                        ("World Health Organization (WHO):", " Global Air Quality Guidelines (2021) — Health thresholds for PM2.5, PM10, O3, NO2, SO2, CO."),
                        ("Central Pollution Control Board (CPCB):", " National Air Quality Index (NAQI) breakpoints and category scales.")
                    ]
                },
                {
                    "heading": "2. Meteorological & Atmospheric Research:",
                    "points": [
                        ("Boundary Layer Meteorology:", " Stull, R. B. (1988) — Planetary Boundary Layer dynamics and thermal inversion mechanics."),
                        ("Atmospheric Chemistry & Physics:", " Seinfeld, J. H., & Pandis, S. N. (2016) — Wet deposition, aerosol growth, and photochemical smog.")
                    ]
                },
                {
                    "heading": "3. Project Repository & Working Prototype:",
                    "points": [
                        ("GitHub Repository:", " https://github.com/NirmalSushil/SIH26082.git"),
                        ("Live Prototype:", " Next.js 16 Coupled Forecast Dashboard with default Mumbai, Favorites/Bookmarks, and Light/Dark Modes.")
                    ]
                }
            ]
        }
    ]

    for slide_data in content_slides_data:
        slide = prs.slides[slide_data["slide_num"]]
        
        # 1. Update and align Team Badge (Oval)
        for shape in slide.shapes:
            if "Oval" in shape.name and shape.has_text_frame:
                shape.left = Inches(0.67)
                shape.top = Inches(0.28)
                shape.width = Inches(2.35)
                shape.height = Inches(0.65)
                
                tf = shape.text_frame
                tf.word_wrap = True
                tf.clear()
                p = tf.paragraphs[0]
                p.text = "Team: Xtream (ID: 09)"
                p.font.size = Pt(11.5)
                p.font.bold = True
                p.font.color.rgb = RGBColor(255, 255, 255)
                p.alignment = PP_ALIGN.CENTER

        # 2. Update and align Slide Title
        for shape in slide.shapes:
            if shape.name == "Title 1" and shape.has_text_frame:
                shape.left = Inches(3.20)
                shape.top = Inches(0.22)
                shape.width = Inches(7.40)
                shape.height = Inches(0.80)
                
                tf = shape.text_frame
                tf.word_wrap = True
                tf.clear()
                p = tf.paragraphs[0]
                p.text = slide_data["title"]
                p.font.size = Pt(17)
                p.font.bold = True
                p.font.color.rgb = COLOR_PRIMARY
                p.alignment = PP_ALIGN.LEFT

        # 3. Update and align Content Textbox
        for shape in slide.shapes:
            if shape.name == "TextBox 8" and shape.has_text_frame:
                shape.left = Inches(0.67)
                shape.top = Inches(1.25)
                shape.width = Inches(11.90)
                shape.height = Inches(5.45)
                
                tf = shape.text_frame
                tf.word_wrap = True
                tf.clear()
                
                for s_idx, sec in enumerate(slide_data["sections"]):
                    p_head = tf.paragraphs[0] if s_idx == 0 else tf.add_paragraph()
                    p_head.space_before = Pt(8) if s_idx > 0 else Pt(0)
                    p_head.space_after = Pt(3)
                    
                    r_head = p_head.add_run()
                    r_head.text = sec["heading"]
                    r_head.font.bold = True
                    r_head.font.size = Pt(15.5)
                    r_head.font.color.rgb = COLOR_PRIMARY
                    
                    for keyword, detail in sec["points"]:
                        p_pt = tf.add_paragraph()
                        p_pt.space_after = Pt(4)
                        p_pt.level = 0
                        
                        # Bullet symbol and keyword
                        r_kw = p_pt.add_run()
                        r_kw.text = "• " + keyword
                        r_kw.font.bold = True
                        r_kw.font.size = Pt(13)
                        r_kw.font.color.rgb = COLOR_HEADER
                        
                        # Detail text
                        r_dt = p_pt.add_run()
                        r_dt.text = detail
                        r_dt.font.bold = False
                        r_dt.font.size = Pt(13)
                        r_dt.font.color.rgb = COLOR_BODY

    prs.save(output_path)
    print(f"Successfully generated SIH presentation with enlarged fonts and aligned layout at: {output_path}")

if __name__ == "__main__":
    generate_perfect_sih_presentation()
